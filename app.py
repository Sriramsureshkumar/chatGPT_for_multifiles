import streamlit as st
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from pptx import Presentation
import openpyxl
import csv
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings, HuggingFaceInstructEmbeddings
from langchain.vectorstores import FAISS
from langchain.chat_models import ChatOpenAI
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from htmlTemplates import css, bot_template, user_template
from langchain.llms import HuggingFaceHub

def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text


from pptx import Presentation

def get_pptx_text(pptx_files):
    text = ""
    for pptx_file in pptx_files:
        presentation = Presentation(pptx_file)
        for slide in presentation.slides:
            slide_text = ""
            
            # Check if the slide has a title
            if slide.shapes.title is not None:
                slide_text += f"Slide Title: {slide.shapes.title.text}\n"

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            slide_text += cell.text + "\t"  # Use a tab to separate cells
            
            text += slide_text
    return text



def get_xlsx_text(xlsx_files):
    text = ""
    for xlsx_file in xlsx_files:
        workbook = openpyxl.load_workbook(xlsx_file)
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            for row in sheet.iter_rows(values_only=True):
                text += " ".join(str(cell) for cell in row) + "\n"
    return text

def get_csv_text(csv_files):
    text = ""
    for csv_file in csv_files:
        with open(csv_file, 'r', encoding='utf-8') as file:
            csv_reader = csv.reader(file)
            for row in csv_reader:
                text += " ".join(row) + "\n"
    return text

def get_text_chunks(text):
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    return chunks


def get_vectorstore(text_chunks):
    embeddings = OpenAIEmbeddings()
    # embeddings = HuggingFaceInstructEmbeddings(model_name="hkunlp/instructor-xl")
    vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
    return vectorstore


def get_conversation_chain(vectorstore):
    llm = ChatOpenAI()
    # llm = HuggingFaceHub(repo_id="google/flan-t5-xxl", model_kwargs={"temperature":0.5, "max_length":512})

    memory = ConversationBufferMemory(
        memory_key='chat_history', return_messages=True)
    conversation_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(),
        memory=memory
    )
    return conversation_chain


def handle_userinput(user_question):
    response = st.session_state.conversation({'question': user_question})
    st.session_state.chat_history = response['chat_history']

    for i, message in enumerate(st.session_state.chat_history):
        if i % 2 == 0:
            st.write(user_template.replace(
                "{{MSG}}", message.content), unsafe_allow_html=True)
        else:
            st.write(bot_template.replace(
                "{{MSG}}", message.content), unsafe_allow_html=True)


def main():
    load_dotenv()
    st.set_page_config(page_title="Chat with multiple documents", page_icon=":books:")
    st.write(css, unsafe_allow_html=True)

    if "conversation" not in st.session_state:
        st.session_state.conversation = None
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = None

    st.header("Chat with multiple documents :books:")
    user_question = st.text_input("Ask a question about your documents:")
    if user_question:
        handle_userinput(user_question)

    with st.sidebar:
        st.subheader("Your documents")
        uploaded_files = st.file_uploader(
            "Upload your documents here and click on 'Process'",
            accept_multiple_files=True,
            type=["pdf", "pptx", "xlsx", "csv"]  # Specify the allowed formats
        )
        if st.button("Process"):
            with st.spinner("Processing"):
                text_chunks = []

                # Process each uploaded file
                for uploaded_file in uploaded_files:
                    if uploaded_file.type == "application/pdf":
                        # Process PDF files
                        raw_text = get_pdf_text([uploaded_file])
                    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                        # Process PPTX files
                        raw_text = get_pptx_text([uploaded_file])
                    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                        # Process XLSX files
                        raw_text = get_xlsx_text([uploaded_file])
                    elif uploaded_file.type == "text/csv":
                        # Process CSV files
                        raw_text = get_csv_text([uploaded_file])

                    # Split and store text chunks as before
                    text_chunks.extend(get_text_chunks(raw_text))

                # Continue with creating a vector store and conversation chain
                vectorstore = get_vectorstore(text_chunks)
                st.session_state.conversation = get_conversation_chain(vectorstore)

if __name__ == '__main__':
    main()
